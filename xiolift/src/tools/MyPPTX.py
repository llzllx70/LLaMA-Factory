
from pptx import Presentation

class MyPPTX:

    @staticmethod
    def pptx_to_text(pptx_path):

        prs = Presentation(pptx_path)
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:  # 检查形状是否包含文本框
                    for paragraph in shape.text_frame.paragraphs:
                        text_content.append(paragraph.text)

        return "\n".join(text_content)